import http.server
import socketserver
import json
import os
import io
import cgi
import traceback

PORT = 5000
HOST = "0.0.0.0"

def extract_text_from_pdf(file_bytes):
    from PyPDF2 import PdfReader
    reader = PdfReader(io.BytesIO(file_bytes))
    text = ""
    for page in reader.pages:
        t = page.extract_text()
        if t:
            text += t + "\n"
    return text

def extract_text_from_docx(file_bytes):
    from docx import Document
    doc = Document(io.BytesIO(file_bytes))
    text = ""
    for para in doc.paragraphs:
        if para.text.strip():
            text += para.text + "\n"
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    text += cell.text + "\n"
    return text

def extract_text_from_pptx(file_bytes):
    from pptx import Presentation
    prs = Presentation(io.BytesIO(file_bytes))
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.strip():
                        text += paragraph.text + "\n"
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            text += cell.text + "\n"
    return text

def build_prompt(text, num_questions, language):
    max_text = text[:12000]
    return f"""You are a medical education expert. Based ONLY on the following educational content, generate exactly {num_questions} multiple-choice questions (MCQs).

CRITICAL RULES:
- Use ONLY the information provided in the content below. Do NOT add external knowledge or information from outside this content.
- Every question, answer, and explanation must be directly based on what is written in the provided content.
- Do NOT use phrases like "According to the content", "Based on the passage", "As mentioned in the text", "According to the provided material" or any similar references to the source content in the questions. Write questions naturally as if they are standalone exam questions.
- Each question must have exactly 4 options (A, B, C, D)
- Only ONE correct answer per question
- Questions should test understanding of the provided material
- Include a brief explanation for the correct answer
- For each WRONG option, include a short explanation of why it is incorrect
- Questions should be in {language}
- Make questions of varying difficulty (easy, medium, hard)

IMPORTANT: Respond ONLY with valid JSON array. No other text before or after.

JSON FORMAT:
[
  {{
    "question": "What is...?",
    "options": ["option A", "option B", "option C", "option D"],
    "correctAnswer": 0,
    "explanation": "Because...",
    "wrongExplanations": {{
      "1": "Why option B is wrong...",
      "2": "Why option C is wrong...",
      "3": "Why option D is wrong..."
    }}
  }}
]

Where correctAnswer is the index (0=A, 1=B, 2=C, 3=D).
wrongExplanations keys are the indices of the wrong options (all options except correctAnswer).

CONTENT:
{max_text}"""

def parse_ai_response(response_text):
    response_text = response_text.strip()
    if response_text.startswith('```'):
        lines = response_text.split('\n')
        lines = lines[1:]
        if lines and lines[-1].strip() == '```':
            lines = lines[:-1]
        response_text = '\n'.join(lines)

    questions = json.loads(response_text)

    if not isinstance(questions, list):
        return {"error": "AI returned invalid format"}

    valid_questions = []
    for q in questions:
        if all(k in q for k in ['question', 'options', 'correctAnswer']):
            if isinstance(q['options'], list) and len(q['options']) == 4:
                vq = {
                    'question': q['question'],
                    'options': q['options'],
                    'correctAnswer': int(q['correctAnswer']),
                    'explanation': q.get('explanation', '')
                }
                if 'wrongExplanations' in q and isinstance(q['wrongExplanations'], dict):
                    vq['wrongExplanations'] = {str(k): str(v) for k, v in q['wrongExplanations'].items()}
                valid_questions.append(vq)

    return {"questions": valid_questions}

# Uses Replit's AI Integrations service for OpenAI-compatible API access
# No API key needed - charges go to Replit credits
def generate_questions_with_openai(text, num_questions, language="English"):
    from openai import OpenAI

    api_key = os.environ.get("AI_INTEGRATIONS_OPENAI_API_KEY")
    base_url = os.environ.get("AI_INTEGRATIONS_OPENAI_BASE_URL")

    print(f"OpenAI check: api_key={'set' if api_key else 'missing'}, base_url={'set' if base_url else 'missing'}")

    if not api_key or not base_url:
        print("OpenAI skipped: missing env vars")
        return None

    client = OpenAI(api_key=api_key, base_url=base_url)
    prompt = build_prompt(text, num_questions, language)

    try:
        # the newest OpenAI model is "gpt-5" which was released August 7, 2025.
        # do not change this unless explicitly requested by the user
        response = client.chat.completions.create(
            model="gpt-5-mini",
            messages=[
                {"role": "system", "content": "You are a medical education expert. Generate MCQ questions and respond ONLY with a JSON object containing a 'questions' array."},
                {"role": "user", "content": prompt}
            ],
            response_format={"type": "json_object"},
            max_completion_tokens=16384
        )
        response_text = response.choices[0].message.content or ""
        print(f"OpenAI response length: {len(response_text)}")
        print(f"OpenAI response preview: {response_text[:300]}")

        try:
            parsed = json.loads(response_text)
        except json.JSONDecodeError:
            return parse_ai_response(response_text)

        questions_list = None
        if isinstance(parsed, dict):
            for key in parsed:
                val = parsed[key]
                if isinstance(val, list) and len(val) > 0:
                    questions_list = val
                    break
        elif isinstance(parsed, list):
            questions_list = parsed

        if questions_list and len(questions_list) > 0:
            valid_questions = []
            for q in questions_list:
                if isinstance(q, dict) and 'question' in q and 'options' in q:
                    opts = q['options']
                    if isinstance(opts, list) and len(opts) == 4:
                        vq = {
                            'question': q['question'],
                            'options': opts,
                            'correctAnswer': int(q.get('correctAnswer', q.get('correct_answer', q.get('answer', 0)))),
                            'explanation': q.get('explanation', '')
                        }
                        if 'wrongExplanations' in q and isinstance(q['wrongExplanations'], dict):
                            vq['wrongExplanations'] = {str(k): str(v) for k, v in q['wrongExplanations'].items()}
                        valid_questions.append(vq)
            if valid_questions:
                has_wrong_exp = sum(1 for q in valid_questions if 'wrongExplanations' in q)
                print(f"OpenAI generated {len(valid_questions)} valid questions ({has_wrong_exp} with wrongExplanations)")
                return {"questions": valid_questions}

        print(f"OpenAI: no valid questions found in response")
        return None
    except Exception as e:
        print(f"OpenAI failed: {type(e).__name__}: {e}")
        return None

def generate_questions_with_gemini(text, num_questions, language="English"):
    import google.generativeai as genai

    api_key = os.environ.get('GOOGLE_API_KEY') or os.environ.get('GEMINI_API_KEY')
    if not api_key:
        return None

    genai.configure(api_key=api_key)
    prompt = build_prompt(text, num_questions, language)

    models_to_try = ['gemini-1.5-flash', 'gemini-2.0-flash-lite', 'gemini-2.0-flash']
    for model_name in models_to_try:
        try:
            model = genai.GenerativeModel(model_name)
            response = model.generate_content(prompt)
            response_text = response.text.strip()
            return parse_ai_response(response_text)
        except Exception as e:
            print(f"Gemini {model_name} failed: {e}")
            continue
    return None

def generate_questions(text, num_questions, language="English"):
    if num_questions > 15:
        all_questions = []
        remaining = num_questions
        batch_num = 0
        while remaining > 0:
            batch_size = min(15, remaining)
            batch_num += 1
            print(f"Batch {batch_num}: generating {batch_size} questions ({remaining} remaining)")
            result = generate_questions_with_openai(text, batch_size, language)
            if result is None:
                result = generate_questions_with_gemini(text, batch_size, language)
            if result and 'questions' in result:
                all_questions.extend(result['questions'])
                remaining -= len(result['questions'])
                if len(result['questions']) < batch_size:
                    print(f"Batch {batch_num} returned fewer than requested, stopping")
                    break
            else:
                print(f"Batch {batch_num} failed")
                break
        if all_questions:
            return {"questions": all_questions}
        return {"error": "Failed to generate questions"}

    result = generate_questions_with_openai(text, num_questions, language)
    if result and 'questions' in result and len(result['questions']) > 0:
        print("Questions generated via OpenAI")
        return result

    result = generate_questions_with_gemini(text, num_questions, language)
    if result and 'questions' in result and len(result['questions']) > 0:
        print("Questions generated via Gemini")
        return result

    return {"error": "Could not generate questions. All AI providers failed. Please try again in a minute."}


class AppHandler(http.server.SimpleHTTPRequestHandler):
    def end_headers(self):
        self.send_header("Cache-Control", "no-cache, no-store, must-revalidate")
        self.send_header("Pragma", "no-cache")
        self.send_header("Expires", "0")
        self.send_header("Access-Control-Allow-Origin", "*")
        self.send_header("Access-Control-Allow-Methods", "POST, GET, OPTIONS")
        self.send_header("Access-Control-Allow-Headers", "Content-Type")
        super().end_headers()

    def do_OPTIONS(self):
        self.send_response(200)
        self.end_headers()

    def do_POST(self):
        if self.path == '/api/generate-questions':
            try:
                content_type = self.headers.get('Content-Type', '')

                if 'multipart/form-data' in content_type:
                    form = cgi.FieldStorage(
                        fp=self.rfile,
                        headers=self.headers,
                        environ={'REQUEST_METHOD': 'POST', 'CONTENT_TYPE': content_type}
                    )

                    if 'file' not in form:
                        self.send_json(400, {"error": "No file uploaded"})
                        return

                    file_item = form['file']
                    if not hasattr(file_item, 'filename') or not file_item.filename:
                        self.send_json(400, {"error": "Invalid file"})
                        return

                    num_questions = int(form.getvalue('numQuestions', '10'))
                    language = form.getvalue('language', 'English')
                    filename = file_item.filename.lower()
                    file_bytes = file_item.file.read()

                    if filename.endswith('.pdf'):
                        text = extract_text_from_pdf(file_bytes)
                    elif filename.endswith('.docx'):
                        text = extract_text_from_docx(file_bytes)
                    elif filename.endswith('.pptx'):
                        text = extract_text_from_pptx(file_bytes)
                    else:
                        self.send_json(400, {"error": "Unsupported file type. Use PDF, Word (.docx), or PowerPoint (.pptx)"})
                        return

                    if not text.strip():
                        self.send_json(400, {"error": "Could not extract text from file. The file might be image-based or empty."})
                        return

                    result = generate_questions(text, num_questions, language)
                    self.send_json(200, result)
                else:
                    self.send_json(400, {"error": "Use multipart/form-data"})
            except Exception as e:
                traceback.print_exc()
                self.send_json(500, {"error": str(e)})
        else:
            self.send_json(404, {"error": "Not found"})

    def send_json(self, code, data):
        self.send_response(code)
        self.send_header('Content-Type', 'application/json')
        self.end_headers()
        self.wfile.write(json.dumps(data, ensure_ascii=False).encode('utf-8'))

socketserver.TCPServer.allow_reuse_address = True
with socketserver.TCPServer((HOST, PORT), AppHandler) as httpd:
    print(f"Serving on http://{HOST}:{PORT}")
    httpd.serve_forever()
